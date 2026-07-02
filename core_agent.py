import os
import json
import io
import operator
from typing import TypedDict, List, Annotated, Literal

import streamlit as st
from langchain_openai import ChatOpenAI
from langgraph.graph import StateGraph, END
from langgraph.checkpoint.memory import MemorySaver


# =========================================================
# 1. 多格式文件解析函数
# =========================================================
def extract_text_from_file(uploaded_file):
    """根据文件类型提取文本：支持 PDF / DOCX / XLSX / PPTX / TXT"""
    file_type = uploaded_file.type
    raw_bytes = uploaded_file.getvalue()

    try:
        if file_type == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            return "".join(page.extract_text() or "" for page in reader.pages)

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(io.BytesIO(raw_bytes))
            return "\n".join(p.text for p in doc.paragraphs)

        elif file_type == "text/plain":
            return raw_bytes.decode("utf-8", errors="ignore")

        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(c) for c in row if c)
                    if row_text:
                        text += row_text + "\n"
            return text

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        else:
            st.error(f"不支持的文件类型: {file_type}")
            return None
    except Exception as e:
        st.error(f"文件解析失败: {e}")
        return None
    
# =========================================================
# 2. 长文本分块 & 去重工具
# =========================================================
def chunk_contract_text(text: str, chunk_size: int = 2500, overlap: int = 200):
    """将长文本按段落切分成固定大小的块，块之间带重叠窗口。

    重叠窗口的作用：防止某个条款恰好卡在两块的分界处被截断。
    例如：'第八条 违约责任：...' 如果刚好被切开，前半段的条款名
    会出现在上一块的末尾（overlap 区域），下一块会重新覆盖到。
    """
    if len(text) <= chunk_size:
        return [text]

    paragraphs = text.split('\n')
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        if len(current_chunk) + len(para) <= chunk_size:
            current_chunk += para + '\n'
        else:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            # 下一块开头带上上一块的尾巴（overlap 字符），防止条款被截断
            tail = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
            current_chunk = tail + para + '\n'

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


def merge_and_deduplicate_clauses(all_clauses: list) -> list:
    """合并去重：同名条款只保留内容最长的那个（信息量最大）"""
    merged = {}
    for clause in all_clauses:
        name = clause.get("clause_name", "")
        content = clause.get("content", "")
        if name not in merged or len(content) > len(merged[name]["content"]):
            merged[name] = clause
    return list(merged.values())


# =========================================================
# 3. 核心合同审核 Agent（基于 LangGraph）
# =========================================================
def run_agent(contract_text: str, api_key: str = None):
    """执行合同审核流程，返回最终状态字典"""
    # 如果没有传入 api_key，则从 st.secrets 读取（云端部署时使用）
    if api_key is None:
        import streamlit as st
        api_key = st.secrets["DEEPSEEK_API_KEY"]
    
    os.environ["OPENAI_API_KEY"] = api_key
    os.environ["OPENAI_API_BASE"] = "https://api.deepseek.com/v1"
    llm = ChatOpenAI(model="deepseek-chat", temperature=0)

    class ContractState(TypedDict):
        contract_text: str
        extracted_clauses: Annotated[List[dict], operator.add]
        risk_report: List[dict]
        current_step: str
        need_human_review: bool

    # ---------- 提取条款 ----------
    def extract_key_clauses(text: str):
        # 短合同直接处理，长合同先分块再逐块提取
        chunks = chunk_contract_text(text)
        all_clauses = []

        for i, chunk in enumerate(chunks):
            chunk_label = f"（第{i+1}/{len(chunks)}段）" if len(chunks) > 1 else ""
            prompt = f"""
        从以下合同文本{chunk_label}中提取关键条款，严格按照JSON列表格式返回。
        请仔细查找合同中与以下维度相关的条款：违约责任、赔偿上限、保密义务、合同终止。
        每条条款提取 clause_name（条款名称）和 content（条款原文摘要）。
        如果实在找不到以上任一维度的条款，返回空列表 []。

        合同文本{chunk_label}：
        {chunk}

        输出示例：
        [{{"clause_name": "违约责任", "content": "违约金为合同总额的20%"}}]
        """
            resp = llm.invoke(prompt).content.strip()
            try:
                if "```" in resp:
                    resp = resp.split("```")[1]
                    if resp.startswith("json"):
                        resp = resp[4:]
                clauses = json.loads(resp)
                if isinstance(clauses, list):
                    all_clauses.extend(clauses)
            except Exception:
                # 单块解析失败不中断全局，继续处理下一块
                continue

        if not all_clauses:
            return [{"clause_name": "解析异常", "content": "所有分块均未能提取到条款，请检查合同格式"}]

        # 同名条款去重：保留内容最完整的版本
        return merge_and_deduplicate_clauses(all_clauses)

    # ---------- 识别风险 ----------
    def identify_risks(clauses: List[dict]):
        if not clauses:
            return []
        prompt = f"""
        请扮演资深法务专家，分析以下合同条款的风险点。
        
        风险等级说明：
        - 高风险：可能导致重大经济损失、法律诉讼或合同无效
        - 中风险：存在法律模糊地带，可能需要谈判调整
        - 低风险：轻微瑕疵，不影响核心利益
        
        必须返回JSON列表格式，每个元素包含：
        - "risk": 风险描述（一句话概括）
        - "level": "高"/"中"/"低"
        - "suggestion": 具体修改建议
        
        如果没有风险，返回空列表 []。
        
        待分析的条款：
        {json.dumps(clauses, ensure_ascii=False)}
        
        只返回JSON列表，不要其他内容。
        """
        resp = llm.invoke(prompt).content.strip()
        try:
            if "```" in resp:
                resp = resp.split("```")[1]
                if resp.startswith("json"):
                    resp = resp[4:]
            return json.loads(resp)
        except:
            return [{"risk": "解析失败", "level": "中", "suggestion": "请人工复核"}]

    # ---------- 节点函数 ----------
    def supervisor_node(state: ContractState):
        if "采购" in state["contract_text"]:
            return {"current_step": "识别为：采购合同"}
        elif "保密" in state["contract_text"]:
            return {"current_step": "识别为：保密协议"}
        return {"current_step": "识别为：通用合同"}

    def extraction_node(state: ContractState):
        st.session_state.progress_text = "正在提取关键条款..."
        return {"extracted_clauses": extract_key_clauses(state["contract_text"])}

    def risk_node(state: ContractState):
        st.session_state.progress_text = "正在分析风险..."
        clauses = state["extracted_clauses"]

        # 条款提取失败时跳过风险分析，直接标记为需人工复核
        if isinstance(clauses, list) and len(clauses) == 1 and clauses[0].get("clause_name") == "解析异常":
            return {
                "risk_report": [{"risk": clauses[0]["content"], "level": "中", "suggestion": "请检查合同文件格式是否正确，或尝试手动粘贴合同内容"}],
                "need_human_review": True
            }

        risks = identify_risks(clauses)
        need = any(r.get("level") == "高" for r in risks) if isinstance(risks, list) else False
        return {"risk_report": risks, "need_human_review": need}

    def human_review_node(state: ContractState):
        st.session_state.progress_text = "⚠️ 触发人工复核（模拟通过）"
        return {"current_step": "人工复核通过"}

    # ---------- 构建图 ----------
    builder = StateGraph(ContractState)
    builder.add_node("supervisor", supervisor_node)
    builder.add_node("extract", extraction_node)
    builder.add_node("risk", risk_node)
    builder.add_node("human_review", human_review_node)

    builder.set_entry_point("supervisor")
    builder.add_edge("supervisor", "extract")
    builder.add_edge("extract", "risk")

    def should_continue(state: ContractState) -> Literal["human_review", "__end__"]:
        return "human_review" if state.get("need_human_review") else END

    builder.add_conditional_edges("risk", should_continue)
    builder.add_edge("human_review", END)

    graph = builder.compile(checkpointer=MemorySaver())

    # ---------- 执行 ----------
    initial_state = {
        "contract_text": contract_text,
        "extracted_clauses": [],
        "risk_report": [],
        "current_step": "",
        "need_human_review": False
    }
    config = {"configurable": {"thread_id": "streamlit_demo"}}
    return graph.invoke(initial_state, config)